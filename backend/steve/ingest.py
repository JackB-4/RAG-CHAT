from typing import List, Tuple
import io
import re
import json
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse

from pypdf import PdfReader
from pdfminer.high_level import extract_text as pdfminer_extract
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def chunk_text(text: str, max_tokens: int = 512, overlap: int = 64) -> List[str]:
    # naive whitespace tokenization for offline use
    words = text.split()
    chunks = []
    step = max_tokens - overlap if max_tokens > overlap else max_tokens
    for i in range(0, len(words), step):
        chunk = words[i : i + max_tokens]
        if chunk:
            chunks.append(" ".join(chunk))
    return chunks


def read_pdf(file_bytes: bytes) -> str:
    # Try PyPDF first
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        text = "\n".join(texts)
        if text and text.strip():
            return clean_text(text)
    except Exception:
        pass
    # Fallback to pdfminer.six
    try:
        text = pdfminer_extract(io.BytesIO(file_bytes))
        return clean_text(text)
    except Exception:
        return ""


def read_docx(file_bytes: bytes) -> str:
    bio = io.BytesIO(file_bytes)
    doc = DocxDocument(bio)
    texts = [p.text for p in doc.paragraphs]
    # include table texts
    for table in doc.tables:
        for row in table.rows:
            texts.append("\t".join([cell.text for cell in row.cells]))
    return clean_text("\n".join(texts))


def read_xlsx(file_bytes: bytes) -> str:
    bio = io.BytesIO(file_bytes)
    wb = load_workbook(bio, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            texts.append("\t".join([str(c) if c is not None else "" for c in row]))
    return clean_text("\n".join(texts))


def read_pptx(file_bytes: bytes) -> str:
    bio = io.BytesIO(file_bytes)
    prs = Presentation(bio)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return clean_text("\n".join(texts))

def read_csv(file_bytes: bytes) -> str:
        # decode bytes as utf-8 with fallback
        try:
            raw = file_bytes.decode('utf-8')
        except Exception:
            raw = file_bytes.decode(errors='ignore')
        # Normalize line endings and compress whitespace
        lines = [clean_text(line) for line in raw.splitlines()]
        return "\n".join(lines)


def read_url(url: str) -> Tuple[str, str]:
    resp = requests.get(url, timeout=30)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "lxml")
    title = soup.title.string.strip() if soup.title and soup.title.string else url
    # remove script/style
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator=" ")
    return title, clean_text(text)
